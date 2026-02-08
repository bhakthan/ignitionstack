"""Stage 2 — Parse: extract structured requirements from input files."""

from __future__ import annotations

import json
from pathlib import Path

from ignition.config import IgnitionConfig
from ignition.llm import chat_json, get_client
from ignition.models import InputType, ParsedRequirements

# ---------------------------------------------------------------------------
# Text extractors (per input type)
# ---------------------------------------------------------------------------


def _extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _extract_pdf(path: Path) -> str:
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except ImportError:
        raise RuntimeError("PyPDF2 is required to parse PDF files. pip install PyPDF2")


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except ImportError:
        raise RuntimeError("python-pptx is required to parse PPTX files. pip install python-pptx")


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document

        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        raise RuntimeError("python-docx is required to parse DOCX files. pip install python-docx")


def _extract_image(path: Path) -> str:
    """Use Vision API to describe a screenshot / image."""
    # Returns base64 placeholder — actual vision call happens in parse()
    import base64

    data = base64.b64encode(path.read_bytes()).decode("utf-8")
    return f"[IMAGE:base64:{data[:50]}...]"  # truncated — full used in vision call


EXTRACTORS = {
    InputType.TEXT: _extract_text,
    InputType.PDF: _extract_pdf,
    InputType.PPTX: _extract_pptx,
    InputType.DOCX: _extract_docx,
    InputType.IMAGE: _extract_image,
}


# ---------------------------------------------------------------------------
# Parse stage
# ---------------------------------------------------------------------------

PARSE_SYSTEM = """\
You are a requirements analyst. Given raw text describing a software project or use case,
extract structured requirements as JSON with these fields:
{
  "summary": "One-paragraph project description",
  "features": ["Feature 1", "Feature 2", ...],
  "constraints": ["Constraint 1", ...],
  "domain_hint": "healthcare | finance | education | oil-and-gas | general",
  "actors": ["Admin", "Patient", ...]
}
Be thorough but concise. Identify all distinct features and constraints."""


def extract_raw_text(path: Path, input_type: InputType) -> str:
    """Extract raw text from the input file."""
    extractor = EXTRACTORS.get(input_type)
    if not extractor:
        raise ValueError(f"Unsupported input type: {input_type}")
    return extractor(path)


def parse(path: Path, input_type: InputType, config: IgnitionConfig) -> ParsedRequirements:
    """Parse an input file into structured requirements using LLM."""
    raw = extract_raw_text(path, input_type)

    # For images, use the vision model
    if input_type == InputType.IMAGE:
        return _parse_image(path, config)

    client = get_client(config)
    result = chat_json(
        client,
        model=config.model,
        system=PARSE_SYSTEM,
        user=f"Extract requirements from this text:\n\n{raw}",
    )
    data = json.loads(result)
    return ParsedRequirements(
        raw_text=raw,
        summary=data.get("summary", ""),
        features=data.get("features", []),
        constraints=data.get("constraints", []),
        domain_hint=data.get("domain_hint", "general"),
        actors=data.get("actors", []),
    )


def _parse_image(path: Path, config: IgnitionConfig) -> ParsedRequirements:
    """Use Vision API for image inputs."""
    import base64

    client = get_client(config)
    b64 = base64.b64encode(path.read_bytes()).decode("utf-8")
    mime = "image/png" if path.suffix.lower() == ".png" else "image/jpeg"

    resp = client.chat.completions.create(
        model=config.model,
        messages=[
            {"role": "system", "content": PARSE_SYSTEM},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract requirements from this screenshot:"},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime};base64,{b64}"},
                    },
                ],
            },
        ],
        response_format={"type": "json_object"},
        max_tokens=4096,
    )
    data = json.loads(resp.choices[0].message.content or "{}")
    return ParsedRequirements(
        raw_text=f"[Image: {path.name}]",
        summary=data.get("summary", ""),
        features=data.get("features", []),
        constraints=data.get("constraints", []),
        domain_hint=data.get("domain_hint", "general"),
        actors=data.get("actors", []),
    )
